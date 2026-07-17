# logic.py
# Imports
import os
import re
import csv
import cv2
import pptx
import fitz
import subprocess
import numpy as np
import numpy.lib.stride_tricks as lst
from mlp import nn, tok_word
from rapidocr_onnxruntime import RapidOCR
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Helper Functions
def grab_text_pptx(temp_path):
    # Opening the Slides
    prs = pptx.Presentation(temp_path)
    slides = prs.slides

    # Going Through All the Slides
    text = []
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text.append(''.join(run.text for run in para.runs))

    return text

def grab_images_pptx(temp_path):
    # Global Helper Variables
    counter = 0

    # Opening the Slides
    prs = pptx.Presentation(temp_path)
    slides = prs.slides

    # Going Through All the Slides
    images = []
    for slide in slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Saving the Image
                im = shape.image
                im_byte = im.blob
                counter += 1

                # Trying to keep it in running memory instead of writing it
                try:
                    img = None
                    if 'wmf' in im.content_type or 'emf' in im.content_type: # Old Slide have this
                        # Using Inkscape to process old image types
                        process = subprocess.run(
                            ["inkscape", "--headless", "--pipe", "--export-type=png", "-o", "-"],
                            input=im_byte,
                            capture_output=True,
                            check=True
                        )
                        arr = np.frombuffer(process.stdout, np.uint8)
                        img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
                    else:
                        # Saving Normal Images
                        arr = np.frombuffer(im_byte, np.uint8)
                        img = cv2.imdecode(arr, cv2.IMREAD_COLOR)

                    if img is not None:
                        images.append(img)
                except Exception as e:
                    print(f"Error -1: Got this error: {e}")
    # Returning the Images
    return images

def grab_text_pdf(temp_path):
    # Reading the PDF
    reader = fitz.open(temp_path)

    # Extracting All Text
    texts = []
    for page in reader:
        texts.append(page.get_text())

    # Returning textz
    return texts

def grab_images_pdf(temp_path):
    # Opening the PDF
    reader = fitz.open(temp_path)

    # Grabbing the Images
    images = []
    for i in range(len(reader)):
        # Loading Page
        page = reader.load_page(i)
        image_list = page.get_images(full = True)

        # Grabbing All the Images
        for img_info in image_list:
            # Grabbing the Raw Image Data
            xref = img_info[0]
            base_image = reader.extract_image(xref)
            im_byte = base_image['image']

            # Adding it to the pathway
            arr = np.frombuffer(im_byte, np.uint8)
            img = cv2.imdecode(arr, cv2.IMREAD_COLOR)
            images.append(img)

    return images

def cloze_function(sentence):
    # Models' HyperParameters
    MAX_WORD_LEN = 45
    MAX_SENT_LEN = 200
    N_EMBD = 100
    WIN_SIZE = 5

    # Opening Model's Parameters
    current_dir = os.path.dirname(__file__)
    data = np.load(os.path.join(current_dir, 'NL-SubwordTaggerCNN.npz'))

    # Extract weights in exact order
    num_params = sum(1 for key in data.files if key.startswith('param_'))
    parameters = [data[f'param_{i}'] for i in range(num_params)]
    C, W1, B1, g1, b1, W2, B2, g2, b2, W3, B3 = parameters

    # 4-Gram Subword Vocabulary
    vocab = data['vocab']
    gramvoid = np.dtype((np.void, vocab.dtype.itemsize * 4))
    vocab_void = np.ascontiguousarray(vocab).view(gramvoid).ravel()

    # Define the unk token and getting index
    unk = np.array([1, 1, 1, 1], dtype=vocab.dtype)
    unk_void = np.ascontiguousarray(unk).view(gramvoid).ravel()
    unk_idx = np.searchsorted(vocab_void, unk_void)

    # Search up func
    def lookup(arr, val, unk_idx):
        if len(arr) == 0:
            return np.ones_like(val) * unk_idx[0]
        idx = np.searchsorted(arr, val)
        cidx = np.minimum(idx, len(arr) - 1)
        return np.where(((idx < len(arr)) & (arr[cidx] == val)), idx, unk_idx[0])

    # Model's Base Char Vocab
    chars = ['t', 'a', '≥', 'β', '!', 'e', 'è', 'κ', '~', ',', '\u200c', 'x', 'v', 'Z', 'ö', '̇', '½', '↑', ';', 'd', ']', 'R', 'é', 'Y', 'ζ', "'", '↔', ' ', 'b', 'γ', 'n', '∞', 'J', '−', 'σ', 'Α', 'f', 's', 'r', '9', 'ü', '↓', 'I', 'E', '/', 'η', '•', '”', '>', 'j', '=', 'm', '{', 'S', '0', 'G', '[', '\\', 'H', '\u202f', 'с', '-', '1', '%', 'y', '&', 'N', 'K', 'à', '_', '°', 'p', '3', 'μ', 'P', '6', '?', 'l', 'V', '≤', 'c', '.', 'T', '’', '"', 'D', 'o', '<', 'Q', '(', 'π', '^', 'g', '#', 'w', '7', 'W', '*', 'L', 'C', 'A', 'F', ')', 'q', 'δ', '8', 'i', 'U', 'z', 'k', '→', 'Ψ', 'h', '+', 'O', 'M', '≈', 'º', 'ç', 'B', 'ε', 'α', '4', '×', '±', '“', '5', '}', 'u', '–', ':', 'X', 'Δ', '—', '2']

    # Recreating String to Integer Index
    stoi = {character: i + 2 for i, character in enumerate(chars)}
    stoi['<PAD>'] = 0
    stoi['<UNK>'] = 1

    def tokenize_char(word):
        ans = [stoi.get(char, 1) for char in word]
        if len(ans) < MAX_WORD_LEN:
            ans += [0] * (MAX_WORD_LEN - len(ans))
        return ans

    words = sentence.split()

    # Tokenize and Pad Characters
    enc_senti = list(map(tokenize_char, words))
    if len(enc_senti) < MAX_SENT_LEN:
        padding = [0] * MAX_WORD_LEN
        enc_senti.extend(padding for _ in range(MAX_SENT_LEN - len(enc_senti)))
    enc_senti = np.array(enc_senti, dtype=vocab.dtype)

    # Transform to 4-Grams and Lookup Vocab Indices
    enc_words = enc_senti.reshape(-1, MAX_WORD_LEN)
    grams = lst.sliding_window_view(enc_words, window_shape=4, axis=1).reshape(-1, 4)
    smpl_void = np.ascontiguousarray(grams).view(gramvoid).ravel()

    indices = lookup(vocab_void, smpl_void, unk_idx)
    xbatch = indices.reshape(1, MAX_SENT_LEN, -1)

    # The Forward Pass
    emb = C[xbatch]
    embcat = emb.reshape(xbatch.shape[0], xbatch.shape[1], -1)

    # Layer 1
    padding_word = np.tile(C[0], MAX_WORD_LEN - 3)
    pad_block = np.tile(padding_word, (embcat.shape[0], 2, 1))
    pembcat = np.concatenate([pad_block, embcat, pad_block], axis=1)

    win1 = lst.sliding_window_view(pembcat, window_shape=WIN_SIZE, axis=1)
    win1r = win1.reshape(win1.shape[:-2] + (-1,))
    win1r = np.ascontiguousarray(win1r)

    h1 = win1r @ W1 + B1
    h1_mean = h1.mean(axis=2, keepdims=True)
    h1_var = h1.var(axis=2, keepdims=True)
    preln1 = (h1 - h1_mean) / ((h1_var + 1e-5) ** 0.5)
    ln1 = g1 * preln1 + b1
    rl1 = np.maximum(np.zeros_like(ln1), ln1)

    # Layer 2
    rl1_padded = np.pad(rl1, ((0, 0), (2, 2), (0, 0)), mode='constant', constant_values=0)
    win2 = lst.sliding_window_view(rl1_padded, window_shape=WIN_SIZE, axis=1)
    win2r = win2.reshape(win2.shape[:-2] + (-1,))
    win2r = np.ascontiguousarray(win2r)

    h2 = win2r @ W2 + B2
    h2_mean = h2.mean(axis=2, keepdims=True)
    h2_var = h2.var(axis=2, keepdims=True)
    preln2 = (h2 - h2_mean) / ((h2_var + 1e-5) ** 0.5)
    ln2 = g2 * preln2 + b2
    rl2 = np.maximum(np.zeros_like(ln2), ln2)

    # Layer 3
    h3 = rl2 @ W3 + B3
    h3 = h3.squeeze(axis=-1)
    probs = 1 / (1 + np.e ** -h3)

    # Decode the Output
    ans = ((probs > 0.5) + 0.0).squeeze(axis=0)
    base = np.ones_like(ans)
    base *= 1/base.shape[0]
    confidence = np.sum(((ans-base)**2)/base) # Using X^2 To determine the confidence

    # if confidence < base.shape[0] - 5:
    #     return ''

    sentence_repieced = ""
    for i, word in enumerate(words):
        if ans[i] == 0:
            sentence_repieced += word + ' '
        else:
            sentence_repieced += '{{' + f'c1::{word}' + '}} '
    return sentence_repieced

def process_text(texts):
    # Preprocessing
    text = ' '.join(texts)
    text = text.replace('\n', ' ')
    sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', text)
    i = 0
    while i < len(sentences):
        words = sentences[i].split()
        if len(words) > 200:
            sentences[i] = ' '.join(words[:200])
            sentences.append(' '.join(words[200:]))
        i += 1


    # Running it Through the Model
    cloze_sentences = []
    for sentence in sentences:
        cloze_sent = cloze_function(sentence)

        # If model deemed it worthy add it to the list
        if len(cloze_sent) > 0:
            cloze_sentences.append(cloze_sent)

    return cloze_sentences

def io_image_prune(image):
    # Hyperparameters
    START_SIZE = 256
    BATCH_SIZE = 32
    WIN_SIZE_L1 = (5, 5)
    NUM_FEATURES_L1 = 16
    MAX_POOLING_SHRINK_L1 = 4
    WIN_SIZE_L2 = (5, 5)
    NUM_FEATURES_L2 = 64
    MAX_POOLING_SHRINK_L2 = 4
    WIN_SIZE_L3 = (5, 5)
    NUM_FEATURES_L3 = 128
    MAX_POOLING_SHRINK_L3 = 2
    LEAKY_RELU_CONSTANT = 0.01
    EPOCHS = 44
    NUM_FEATURES_L4 = 64
    LR = 0.04

    # Loading the Parameters
    current_dir = os.path.dirname(__file__)
    data = np.load(os.path.join(current_dir, 'NL-ImageClassifierCNN.npz'))
    W1, g1, b1, W2, g2, b2, W3, g3, b3, W4, g4, b4, W5 = list(data.values())[:13] # Dont ask y I am only indexing the first 12 elements

    # Preprocessing the image
    img = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    resized_img = cv2.resize(img, (START_SIZE, START_SIZE))
    xbatch = np.expand_dims(resized_img, axis=0)

    ### Forward Pass
    ## Layer 1
    # Padding
    pad1 = tuple((size - 1) // 2 for size in WIN_SIZE_L1)  # Same Padding Calcs
    img_padded1 = np.pad(xbatch, ((0, 0), pad1, pad1), "constant", constant_values=0)

    # Conv
    win_l1 = lst.sliding_window_view(img_padded1, window_shape=WIN_SIZE_L1, axis=(1, 2))
    win_l1r = (np.ascontiguousarray(win_l1)).reshape(-1, win_l1.shape[-2] * win_l1.shape[-1])

    # Linear
    preln1 = win_l1r @ W1.reshape(-1, W1.shape[-1])
    preln1 = preln1.reshape(-1, xbatch.shape[1], xbatch.shape[1], NUM_FEATURES_L1)

    # Layer Norm
    lnmean1 = np.mean(preln1, axis=-1, keepdims=True)
    lnvar1 = np.var(preln1, axis=-1, keepdims=True)
    lnraw1 = (preln1 - lnmean1) / np.sqrt(lnvar1 + 1e-5)
    ln1 = g1 * lnraw1 + b1

    # ReLu
    rl1 = np.maximum(LEAKY_RELU_CONSTANT * ln1, ln1)

    # Max Pooling
    rl1t = rl1.transpose(0, 3, 1, 2)
    img_pieced1 = rl1t.reshape(-1, NUM_FEATURES_L1, rl1t.shape[-1] // MAX_POOLING_SHRINK_L1, MAX_POOLING_SHRINK_L1,
                               rl1t.shape[-1] // MAX_POOLING_SHRINK_L1, MAX_POOLING_SHRINK_L1)
    pooled_img1 = np.max(img_pieced1, axis=(-1, -3))

    ## Layer 2
    # Padding
    pad2 = tuple((size - 1) // 2 for size in WIN_SIZE_L2)
    img_padded2 = np.pad(pooled_img1, ((0, 0), (0, 0), pad2, pad2), "constant", constant_values=0)

    # Conv
    win_l2 = lst.sliding_window_view(img_padded2, window_shape=WIN_SIZE_L2, axis=(2, 3))
    win_l2c = np.ascontiguousarray(win_l2).transpose(0, 2, 3, 4, 5, 1)
    win_l2r = win_l2c.reshape(-1, win_l2c.shape[-2] * win_l2c.shape[-1] * win_l2c.shape[-3])

    # Linear
    preln2 = win_l2r @ W2.reshape(-1, W2.shape[-1])
    preln2 = preln2.reshape(-1, pooled_img1.shape[2], pooled_img1.shape[2], NUM_FEATURES_L2)

    # Layer Norm
    lnmean2 = np.mean(preln2, axis=-1, keepdims=True)
    lnvar2 = np.var(preln2, axis=-1, keepdims=True)
    lnraw2 = (preln2 - lnmean2) / np.sqrt(lnvar2 + 1e-5)
    ln2 = g2 * lnraw2 + b2

    # ReLu
    rl2 = np.maximum(LEAKY_RELU_CONSTANT * ln2, ln2)

    # Max Pooling
    rl2t = rl2.transpose(0, 2, 3, 1)
    img_pieced2 = rl2t.reshape(-1, NUM_FEATURES_L2, rl2t.shape[-1] // MAX_POOLING_SHRINK_L2, MAX_POOLING_SHRINK_L2,
                               rl2t.shape[-1] // MAX_POOLING_SHRINK_L2, MAX_POOLING_SHRINK_L2)
    pooled_img2 = np.max(img_pieced2, axis=(-1, -3))

    ## Layer 3
    # Padding
    pad3 = tuple((size - 1) // 2 for size in WIN_SIZE_L3)
    img_padded3 = np.pad(pooled_img2, ((0, 0), (0, 0), pad3, pad3), "constant", constant_values=0)

    # Conv
    win_l3 = lst.sliding_window_view(img_padded3, window_shape=WIN_SIZE_L3, axis=(2, 3))
    win_l3c = np.ascontiguousarray(win_l3).transpose(0, 2, 3, 4, 5, 1)
    win_l3r = win_l3c.reshape(-1, win_l3c.shape[-2] * win_l3c.shape[-1] * win_l3c.shape[-3])

    # Linear
    preln3 = win_l3r @ W3.reshape(-1, W3.shape[-1])
    preln3 = preln3.reshape(-1, pooled_img2.shape[2], pooled_img2.shape[2], NUM_FEATURES_L3)

    # Layer Norm
    lnmean3 = np.mean(preln3, axis=-1, keepdims=True)
    lnvar3 = np.var(preln3, axis=-1, keepdims=True)
    lnraw3 = (preln3 - lnmean3) / np.sqrt(lnvar3 + 1e-5)
    ln3 = g3 * lnraw3 + b3

    # ReLu
    rl3 = np.maximum(LEAKY_RELU_CONSTANT * ln3, ln3)

    # Max Pooling
    rl3t = rl3.transpose(0, 3, 1, 2)
    img_pieced3 = rl3t.reshape(-1, NUM_FEATURES_L3, rl3t.shape[-1] // MAX_POOLING_SHRINK_L3, MAX_POOLING_SHRINK_L3,
                               rl3t.shape[-1] // MAX_POOLING_SHRINK_L3, MAX_POOLING_SHRINK_L3)
    pooled_img3 = np.max(img_pieced3, axis=(-1, -3))

    ## Layer 4
    img_flattened = pooled_img3.reshape(img_pieced3.shape[0], -1)
    h4 = img_flattened @ W4
    # mlp_data = h4.copy() I am removing it. It'll be too hard to train the word dataset unless I give Gemini d images which is slow

    # Layer Norm... sighs
    lnmean4 = h4.mean(-1, keepdims=True)
    lnvar4 = h4.var(-1, keepdims=True)
    lnraw4 = (h4 - lnmean4) / np.sqrt(lnvar4 + 1e-5)
    ln4 = g4 * lnraw4 + b4

    # ReLu
    rl4 = np.maximum(LEAKY_RELU_CONSTANT * ln4, ln4)

    ## Layer 5
    raw_pred = rl4 @ W5

    ## Sigmoid + Loss
    safe_raw_pred = np.clip(raw_pred, -250, 250)
    probs = 1 / (1 + (np.e ** -(safe_raw_pred)))

    if probs[0][0] > 0.5:
        return 1
    return 0

def io_word_prune(word):
    # Tokenize the word
    x_input = np.array([tok_word(word.lower())])

    # Forward pass through MLP
    logits = nn(x_input)

    # Apply Sigmoid
    probs = (1 + np.e ** -logits + 1e-5) ** -1

    # Get Classification
    if probs[0][0] > 0.5:
        return 1
    return 0

def occlude_word(image, cord, color=(0, 255, 0)):
    # Reshaping cords to work w/ cv2
    cord = np.array(cord).astype(np.int32)
    reshaped_cord = cord.reshape(-1, 1, 2)

    # Filling At The Spot
    img = cv2.fillPoly(image, [reshaped_cord], color)
    return img

def io_generate(counter, image, cords, temp_dir):
    for cord in cords:
        image = occlude_word(image, cord)

    for select_cord in cords:
        img = occlude_word(image.copy(), select_cord, (0, 0, 255))
        cv2.imwrite(f'{temp_dir}/image_{counter}.png', img)
        counter += 1

    return counter

def process_images(images, temp_dir):
    # Loading OCR Engine
    ocr_engine = RapidOCR()

    # Scanning the Images
    counter = 0
    words = []
    for image in images:
        results, elapse = ocr_engine(image)

        cords = []
        if results:
            if io_image_prune(image) == 1:
                for detection in results:
                    text = detection[1]
                    confidence = detection[2]
                    if re.search('[a-zA-Z]', text) and confidence > 0.4:
                       if io_word_prune(text) == 1:
                           cords.append(detection[0])
                           words.append(text)

        counter = io_generate(counter, image, cords, temp_dir)

    return words

# Main Functions
def run_pipeline(temp_path, temp_dir):
    try:
        images = None
        texts = None
        if temp_path.lower().endswith('.pptx'):
            # Grabbing the Images and Saving them to Local Temp Dir
            images = grab_images_pptx(temp_path)

            # Grabbing the Text on the Slides
            texts = grab_text_pptx(temp_path)
        elif temp_path.lower().endswith('.pdf'):
            # Same thing but w/ PDF this time
            images = grab_images_pdf(temp_path)

            # PDF style
            texts = grab_text_pdf(temp_path)
        else:
            print(f'ERROR 0: Got {temp_path.lower()} as file expected .pdf or .pptx')

        words, cloze_sentences = None, None
        if images is not None:
            words = process_images(images, temp_dir)

        if texts is not None:
            cloze_sentences = process_text(texts)

        return words, cloze_sentences
    except Exception as e:
        print(f'ERROR -1: {e}')
        raise e


def create_csv_file(deck_name, temp_path, temp_dir):
    # Running the Pipeline
    words, cloze_sentences = run_pipeline(temp_path, temp_dir)

    # Writing the CSV
    csv_filename = f"{deck_name}.csv"
    csv_filepath = os.path.join(temp_dir, csv_filename)
    csv_rows = []

    # Creating IO Cards first
    if words is not None:
        for index, word in enumerate(words):
            front_html = f'<img src="image_{index}.png">'
            back_text = word
            csv_rows.append([front_html, back_text])

    # Creating Cloze Sentences Next
    if cloze_sentences is not None:
        for sentence in cloze_sentences:
            csv_rows.append([sentence, ""])

    # Returning the CSV
    if csv_rows:
        with open(csv_filepath, 'w', newline='', encoding='utf-8') as csvfile:
            writer = csv.writer(csvfile)
            writer.writerows(csv_rows)
        return csv_filename

    return None
